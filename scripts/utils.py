from googleapiclient.http import MediaFileUpload
from pptx import Presentation
import requests
import os

def create_ppt_with_date_and_members(date, save_path, filename, members):
    script_dir = os.path.dirname(os.path.abspath(__file__))
    template_path = os.path.join(script_dir, "Presentation_template.pptx")
    
    prs = Presentation(template_path)
    
    slide = prs.slides[0]
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "5 minutes presentations"
    subtitle.text = f"Date: {date}"
    
    for member in members:
        slide_layout = prs.slide_layouts[2]
        slide = prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        title.text = f"{member} - Previous week"
        
        slide = prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        title.text = f"{member} - This week"
    
    file_path = os.path.join(save_path, filename)

    if os.path.exists(file_path):
        os.remove(file_path)
    
    prs.save(file_path)
    return file_path

def upload_to_drive(service, file_path, filename, folder_id):
    file_metadata = {'name': filename, 'parents': [folder_id]}
    media = MediaFileUpload(file_path, mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation')
    file = service.files().create(body=file_metadata, media_body=media, fields='id').execute()
    return file.get('id')

def create_shareable_link(service, file_id):
    permission = {'type': 'anyone', 'role': 'writer'}
    service.permissions().create(fileId=file_id, body=permission).execute()
    return f"https://drive.google.com/file/d/{file_id}/view"

# Function to send message to Slack
def send_slack_message(token, channel, text):
    url = 'https://slack.com/api/chat.postMessage'
    headers = {
        'Content-Type': 'application/json',
        'Authorization': f'Bearer {token}'
    }
    data = {
        'channel': channel,
        'text': text
    }
    response = requests.post(url, headers=headers, json=data)
    if response.status_code != 200:
        raise Exception(f"Request to Slack API failed with status code {response.status_code}, response: {response.text}")
